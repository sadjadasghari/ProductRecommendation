import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_slide(prs, title, content):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.shapes.placeholders[1]
    
    title_shape.text = title
    
    tf = body_shape.text_frame
    tf.clear()
    
    for i, bullet in enumerate(content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(24)

def create_image_slide(prs, title, image_path):
    slide_layout = prs.slide_layouts[5] # Title only
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # insert image centered
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    
    try:
        slide.shapes.add_picture(image_path, left, top, width)
    except Exception as e:
        print(f"Could not load image {image_path}: {e}")

def generate_presentation():
    # Fix dict attribute error sometimes seen in older pptx versions with python 3.10+
    prs = Presentation()
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Multimodal Intelligent System Experiences (ISE)"
    subtitle.text = "A Privacy-First Architecture for Personalized Product Discovery\nCandidate Presentation: Manager, SIML"
    
    # Slides Content
    slides_data = [
        {
            "title": "1. Vision & Strategy (SIML Focus)",
            "content": [
                "Goal: Deliver highly personalized, multimodal intelligent experiences natively on Apple devices.",
                "System Intelligence: Decoupling heavy representation learning from on-device sequential modeling.",
                "Privacy-First: User click streams remain entirely on-device.",
                "Scalability: Designed to scale across Apple's ecosystem (iOS, iPadOS, macOS) via CoreML and the Apple Neural Engine (ANE)."
            ]
        },
        {
            "title": "2. Multimodal Foundation Models",
            "content": [
                "Catalog Embedding Cloud: Pre-processing the product index purely offline.",
                "Vision Backbone: ViT-B/16 (Vision Transformer) for semantic patch-level feature extraction.",
                "Text Backbone: all-MiniLM Transformer replacing RNNs for deeply contextualized text alignment.",
                "Result: A unified 128-dimensional multimodal item embedding representing complex semantics."
            ]
        },
        {
            "title": "3. System Architecture Design",
            "image": "architecture.png"
        },
        {
            "title": "4. Efficient On-Device ML Architecture",
            "content": [
                "RecSys Edge Tower: A fast SASRec (Self-Attentive) sequence Transformer processing local interaction history.",
                "Hard Negative Contrastive Loss: Trained symmetrically using Additive Margin (CosFace) mechanics.",
                "Feature Store Integration: Injecting real-time context (time-of-day, location) efficiently onto the edge.",
                "Retrieval Engine: Employing FAISS HNSW for sub-millisecond Approximate Nearest Neighbor resolution natively in RAM.",
                "Quantization: Compressing core weights to INT8 to radically minimize runtime footprint."
            ]
        },
        {
            "title": "5. Agentic Generative AI (RAG + Diffusion)",
            "content": [
                "Intelligent Orchestration: Multimodal LLM dynamically routes ambiguous user queries.",
                "Spatial Inpainting: Users upload photos of their own space (e.g., Living Room) for contextualized placement.",
                "Closed-Loop Feedback: User validations of Spatial images implicitly construct difficult hard negatives.",
                "LoRA Fine-Tuning: Applying Low-Rank Adaptation to strictly adhere to brand guidelines and aesthetics."
            ]
        },
        {
            "title": "6. End-to-End Evaluation Framework",
            "content": [
                "Recommendation Quality: Driving engagement via Hit Rate (HR@10) and NDCG via dense indexing.",
                "Generative AI Metrics: Automated CLIP Score evaluation for text-to-image alignment and prompt adherence.",
                "Continuous Improvement: Setting up infrastructure for RLHF and online A/B testing (CTR tracking)."
            ]
        },
        {
            "title": "7. Team Leadership & Roadmap",
            "content": [
                "Cross-Functional Execution: Bridging researchers, mobile developers, and product teams.",
                "Infrastructure: Standardizing on scalable tools (PyTorch Lightning, Ray Serve) vs Edge deployments (CoreML).",
                "Next Steps: Expanding pure LLM integration for conversational product search and Multimodal Foundation capabilities."
            ]
        }
    ]
    
    for slide_data in slides_data:
        if "image" in slide_data:
            create_image_slide(prs, slide_data["title"], slide_data["image"])
        else:
            create_slide(prs, slide_data["title"], slide_data["content"])
        
    prs.save('apple_siml_interview_presentation.pptx')
    print("✅ Successfully generated 'apple_siml_interview_presentation.pptx'")

if __name__ == "__main__":
    generate_presentation()
